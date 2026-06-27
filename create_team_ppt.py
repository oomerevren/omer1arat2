from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pathlib import Path

OUT = Path('/home/user/deeppipelinedsad/reports/Deep-Pipeline_Takim_Tanitim.pptx')
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Theme colors
NAVY = RGBColor(10, 18, 34)
NAVY2 = RGBColor(15, 28, 52)
CYAN = RGBColor(0, 210, 255)
PURPLE = RGBColor(124, 77, 255)
GREEN = RGBColor(45, 212, 191)
ORANGE = RGBColor(255, 184, 77)
WHITE = RGBColor(245, 248, 255)
MUTED = RGBColor(160, 175, 205)
GRAY = RGBColor(52, 66, 92)
RED = RGBColor(255, 91, 91)

FONT_HEAD = 'Aptos Display'
FONT_BODY = 'Aptos'


def add_bg(slide, title=None, subtitle=None, page=None):
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid(); bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()
    # Accent shapes
    c1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.6), Inches(-1.0), Inches(3.4), Inches(3.4))
    c1.fill.solid(); c1.fill.fore_color.rgb = RGBColor(22, 52, 92); c1.fill.transparency = 20
    c1.line.fill.background()
    c2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-1.4), Inches(5.4), Inches(3.3), Inches(3.3))
    c2.fill.solid(); c2.fill.fore_color.rgb = RGBColor(36, 28, 79); c2.fill.transparency = 18
    c2.line.fill.background()
    # Top line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.07))
    line.fill.solid(); line.fill.fore_color.rgb = CYAN
    line.line.fill.background()
    # Header mini logo
    logo = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.45), Inches(0.25), Inches(1.68), Inches(0.42))
    logo.fill.solid(); logo.fill.fore_color.rgb = RGBColor(16, 37, 68)
    logo.line.color.rgb = CYAN; logo.line.width = Pt(1)
    tx = logo.text_frame; tx.clear(); p=tx.paragraphs[0]
    p.text = 'Deep-Pipeline'; p.font.name=FONT_BODY; p.font.size=Pt(12); p.font.bold=True; p.font.color.rgb=WHITE; p.alignment=PP_ALIGN.CENTER
    tx.vertical_anchor = MSO_ANCHOR.MIDDLE
    # Footer
    foot = slide.shapes.add_textbox(Inches(0.45), Inches(7.05), Inches(8), Inches(0.25))
    tf=foot.text_frame; tf.clear(); p=tf.paragraphs[0]
    p.text='TEKNOFEST 2026 E-Ticaret Hackathonu  •  Takım Tanıtım Dosyası'
    p.font.name=FONT_BODY; p.font.size=Pt(9); p.font.color.rgb=MUTED
    if page:
        pg = slide.shapes.add_textbox(Inches(12.25), Inches(7.05), Inches(0.6), Inches(0.25))
        t=pg.text_frame; t.clear(); p=t.paragraphs[0]; p.text=str(page); p.font.name=FONT_BODY; p.font.size=Pt(9); p.font.color.rgb=MUTED; p.alignment=PP_ALIGN.RIGHT
    if title:
        box=slide.shapes.add_textbox(Inches(0.65), Inches(0.86), Inches(12.0), Inches(0.65))
        tf=box.text_frame; tf.clear(); p=tf.paragraphs[0]
        p.text=title; p.font.name=FONT_HEAD; p.font.size=Pt(30); p.font.bold=True; p.font.color.rgb=WHITE
    if subtitle:
        box=slide.shapes.add_textbox(Inches(0.68), Inches(1.48), Inches(11.8), Inches(0.35))
        tf=box.text_frame; tf.clear(); p=tf.paragraphs[0]
        p.text=subtitle; p.font.name=FONT_BODY; p.font.size=Pt(13); p.font.color.rgb=MUTED


def textbox(slide, x, y, w, h, text='', size=16, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    shp=slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf=shp.text_frame; tf.clear(); tf.word_wrap=True
    p=tf.paragraphs[0]; p.text=text; p.font.name=FONT_BODY; p.font.size=Pt(size); p.font.color.rgb=color; p.font.bold=bold; p.alignment=align
    return shp


def pill(slide, x, y, w, h, text, fill, color=WHITE, size=11, line=None):
    shp=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb=fill
    shp.line.color.rgb = line if line else fill
    shp.line.width=Pt(1)
    tf=shp.text_frame; tf.clear(); tf.vertical_anchor=MSO_ANCHOR.MIDDLE
    p=tf.paragraphs[0]; p.text=text; p.font.name=FONT_BODY; p.font.size=Pt(size); p.font.bold=True; p.font.color.rgb=color; p.alignment=PP_ALIGN.CENTER
    return shp


def card(slide, x, y, w, h, title, body, accent=CYAN, title_size=16, body_size=11):
    shp=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb=NAVY2
    shp.line.color.rgb=RGBColor(38, 66, 110); shp.line.width=Pt(1)
    # accent bar
    bar=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.06), Inches(h))
    bar.fill.solid(); bar.fill.fore_color.rgb=accent; bar.line.fill.background()
    tb=slide.shapes.add_textbox(Inches(x+0.18), Inches(y+0.15), Inches(w-0.35), Inches(0.35))
    tf=tb.text_frame; tf.clear(); p=tf.paragraphs[0]; p.text=title; p.font.name=FONT_HEAD; p.font.size=Pt(title_size); p.font.bold=True; p.font.color.rgb=WHITE
    bb=slide.shapes.add_textbox(Inches(x+0.18), Inches(y+0.58), Inches(w-0.35), Inches(h-0.7))
    tf=bb.text_frame; tf.clear(); tf.word_wrap=True
    for i, line in enumerate(body.split('\n')):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.text=line; p.font.name=FONT_BODY; p.font.size=Pt(body_size); p.font.color.rgb=MUTED; p.space_after=Pt(3)
    return shp


def bullets(slide, x, y, w, h, items, size=13, color=MUTED, bullet_color=None):
    tb=slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf=tb.text_frame; tf.clear(); tf.word_wrap=True
    for i, item in enumerate(items):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.text=item; p.font.name=FONT_BODY; p.font.size=Pt(size); p.font.color.rgb=color; p.level=0; p.space_after=Pt(5)
        p._p.get_or_add_pPr().set('marL', '285750')
        p._p.get_or_add_pPr().set('indent', '-171450')
        # simple bullet glyph
        p.text = '• ' + item
    return tb

# Slide 1 cover
slide=prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide, page='01')
textbox(slide, 0.75, 1.15, 11.8, 0.55, 'DEEP-PIPELINE', 40, WHITE, True)
textbox(slide, 0.78, 1.78, 11.4, 0.48, 'E-Ticaret Arama Alaka Düzeyi için Uçtan Uca Açıklanabilir ML Pipeline', 18, GREEN, True)
textbox(slide, 0.80, 2.35, 10.0, 0.7, 'TEKNOFEST 2026 E-Ticaret Hackathonu kapsamında; Türkçe e-ticaret metinlerinde sorgu–ürün ilişkisinin daha doğru, hızlı ve açıklanabilir sınıflandırılması.', 15, MUTED)
pill(slide, 0.80, 3.25, 2.3, 0.42, 'Türkçe NLP', RGBColor(20, 78, 92), size=12, line=GREEN)
pill(slide, 3.25, 3.25, 2.7, 0.42, 'Cross-Encoder', RGBColor(37, 41, 92), size=12, line=PURPLE)
pill(slide, 6.10, 3.25, 2.0, 0.42, 'XAI', RGBColor(20, 78, 92), size=12, line=CYAN)
pill(slide, 8.25, 3.25, 2.9, 0.42, 'MLOps + Docker', RGBColor(58, 48, 28), size=12, line=ORANGE)
# Big pipeline visual
for i,(label, col) in enumerate([('Veri', CYAN), ('Model', PURPLE), ('API', GREEN), ('Sunum', ORANGE)]):
    x=1.0+i*2.65
    card(slide, x, 4.25, 2.1, 1.05, label, ['EDA & normalize','CE + ensemble','FastAPI + XAI','Jüri demosu'][i], col, 17, 11)
    if i<3:
        conn=slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x+2.12), Inches(4.78), Inches(x+2.62), Inches(4.78))
        conn.line.color.rgb=MUTED; conn.line.width=Pt(2)
textbox(slide, 0.80, 6.25, 9.5, 0.28, 'Takım: Ömer Mahmut Evren • Ahmet Arat Turmuş • Danışman: Mustafa Gökmen', 13, WHITE, True)
textbox(slide, 0.80, 6.58, 6.5, 0.25, 'Tarih: 18 Haziran 2026', 10, MUTED)

# Slide 2 Problem
slide=prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide, 'Problem ve Vizyon', 'E-ticarette kullanıcı sorgusu ile ürün içeriği arasındaki anlam eşleşmesini güçlendirmek.', '02')
card(slide, 0.75, 2.0, 3.75, 3.0, 'Yarışma Problemi', 'Kullanıcı sorgusu ile ürün başlığı / marka / kategori / özellikler arasında alaka düzeyini sınıflandırmak.\n\nHedef: “Alakasız”, “Kısmen Alakalı”, “Çok Alakalı” kararını tutarlı üretmek.', CYAN, 18, 12)
card(slide, 4.8, 2.0, 3.75, 3.0, 'Neden Zor?', '• Türkçe ekler, yazım hataları ve kısaltmalar\n• Marka–kategori çakışmaları\n• Benzer başlıklı ama farklı niyetli ürünler\n• Macro-F1 dengesinde sınıf ayrımı', PURPLE, 18, 12)
card(slide, 8.85, 2.0, 3.75, 3.0, 'Bizim Vizyonumuz', 'Sadece yüksek skor alan değil; jüriye nasıl karar verdiğini gösterebilen, tekrarlanabilir ve üretime taşınabilir bir sistem.', GREEN, 18, 12)
textbox(slide, 1.0, 5.65, 11.2, 0.5, 'Deep-Pipeline: Kaggle performansı + açıklanabilir final demo + mühendislik kanıtı üçlüsünü aynı dosya yapısında birleştirir.', 18, WHITE, True, PP_ALIGN.CENTER)

# Slide 3 Solution
slide=prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide, 'Çözüm Özeti', 'Veriden sunuma kadar tek giriş noktası olan profesyonel bir ML sistemi.', '03')
# pipeline blocks
steps=[('01\nVeri', 'train/test hazırlama\nnormalize labels', CYAN), ('02\nÖn İşleme', 'Türkçe lower\nSymSpell + Zeyrek', GREEN), ('03\nÖzellik', 'BM25, marka, kategori\nfuzzy matching', ORANGE), ('04\nModel', 'DistilBERTurk CE\nensemble + threshold', PURPLE), ('05\nServis', 'FastAPI /predict\n/explain /metrics', CYAN), ('06\nSunum', 'XAI paneli\n3D jüri demosu', GREEN)]
for i,(t,b,c) in enumerate(steps):
    x=0.6+i*2.1
    card(slide, x, 2.05, 1.78, 2.1, t, b, c, 16, 9.5)
    if i<5:
        conn=slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x+1.80), Inches(3.1), Inches(x+2.08), Inches(3.1))
        conn.line.color.rgb=MUTED; conn.line.width=Pt(1.7)
# bottom details
card(slide, 0.85, 4.75, 3.65, 1.45, 'Skor Stratejisi', 'Macro-F1 odaklı validation, kategori bazlı eşik arama, hard negative mining ve veri sonrası hızlı deney döngüsü.', CYAN, 16, 11)
card(slide, 4.85, 4.75, 3.65, 1.45, 'Açıklanabilirlik', 'Feature katkısı + transformer attention görselleştirmesi ile karar gerekçesi sunulur.', GREEN, 16, 11)
card(slide, 8.85, 4.75, 3.65, 1.45, 'Dağıtım', 'Docker, FastAPI ve offline model yolu ile final ortamına taşınabilir servis mimarisi.', ORANGE, 16, 11)

# Slide 4 Differentiators
slide=prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide, 'Özgünlük ve Güçlü Yanlar', 'Deep-Pipeline yalnızca model değil; jüriye gösterilebilir mühendislik sistemi.', '04')
items=[
    ('Türkçe E-Ticaret NLP', 'İ/ı dönüşümü, yazım hatası düzeltme, morfolojik analiz ve e-ticaret sözlüğü.'),
    ('Hard Negative Mining', 'Modelin en çok karıştırdığı yakın ama yanlış ürünleri eğitim sinyaline dönüştürme.'),
    ('Cross-Encoder + Ensemble', 'Sorgu–ürün çiftini birlikte okuyan semantik model; özellik tabanlı sinyallerle güçlendirme.'),
    ('Açıklanabilir AI', 'Jüri demosunda “neden bu karar?” sorusuna feature ve attention katmanı ile cevap.'),
    ('MLOps Disiplini', 'YAML config, MLflow deney takibi, testler, Docker ve yeniden üretilebilir komutlar.'),
    ('Sunum Stratejisi', 'Teknik doğruluk, kullanıcı hikâyesi ve canlı demo aynı anlatı içinde.'),
]
for idx,(t,b) in enumerate(items):
    row=idx//2; col=idx%2
    x=0.75+col*6.1; y=1.9+row*1.55
    card(slide, x, y, 5.65, 1.22, t, b, [CYAN,GREEN,PURPLE,ORANGE,CYAN,GREEN][idx], 15, 10.2)

# Slide 5 Current Status
slide=prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide, 'Mevcut Durum ve Kanıtlar', 'Başvuru tamamlandı; teknik altyapı veri seti açıklanınca çalıştırılmaya hazır.', '05')
# status table cards
card(slide, 0.75, 1.85, 3.85, 1.1, 'KYS Başvuru', 'Kayıt, takım kurulumu ve başvuru tamamlandı. Google Group takibi başladı.', GREEN, 16, 11)
card(slide, 4.75, 1.85, 3.85, 1.1, 'Repo Altyapısı', 'src/ tabanlı modüler yapı; data, training, evaluation, deployment ve xai modülleri.', CYAN, 16, 11)
card(slide, 8.75, 1.85, 3.85, 1.1, 'Test / API', 'FastAPI uçları, label testleri ve pipeline komutları ile doğrulanabilir yapı.', PURPLE, 16, 11)
# timeline
textbox(slide, 0.8, 3.55, 11.2, 0.4, 'Yakın Yol Haritası', 20, WHITE, True)
for i,(date,title,body,col) in enumerate([
    ('18–19 Haz', 'Aşama 1', 'Takım tanıtımı, görev netliği, KYS dosyaları', GREEN),
    ('26 Haz', 'Veri Açılışı', 'Kaggle train/test hazırlama ve baseline', CYAN),
    ('Temmuz', 'Skor Artırımı', 'CV, threshold, ensemble, ablation', PURPLE),
    ('Final', 'Jüri Demosu', 'XAI, latency, rapor ve canlı sunum', ORANGE),
]):
    x=0.95+i*3.05
    pill(slide, x, 4.18, 1.15, 0.38, date, col, NAVY, 10)
    card(slide, x, 4.70, 2.55, 1.35, title, body, col, 14, 9.7)
    if i<3:
        conn=slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x+2.58), Inches(5.38), Inches(x+3.0), Inches(5.38))
        conn.line.color.rgb=MUTED; conn.line.width=Pt(2)
textbox(slide, 0.8, 6.42, 11.7, 0.28, 'Not: Gerçek yarışma verisi açıklanana kadar performans metrikleri iddia değil, veri sonrası ölçülecek hedef olarak ele alınır.', 10.5, MUTED)

# Slide 6 Team roles
slide=prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide, 'Takım ve Rol Dağılımı', 'Öğrenci liderliğinde, danışman gözetimli ve görevleri izlenebilir bir çalışma modeli.', '06')
# leader big card
card(slide, 0.75, 1.9, 5.65, 3.95, 'Ömer Mahmut Evren\nTakım Kaptanı • ML/System Lead', '• Proje vizyonu ve teknik kararların ana sorumlusu\n• Model mimarisi, training pipeline ve Kaggle stratejisi\n• Repo mimarisi, deney planı, rapor ve jüri anlatısının sahibi\n• Final sunumunda teknik akış ve demo liderliği', CYAN, 20, 12)
# supporting cards
card(slide, 6.75, 1.9, 2.85, 3.95, 'Ahmet Arat Turmuş\nQA & Dokümantasyon Destek', '• Başvuru ve dosya kontrol checklist’i\n• Çalıştırma adımlarını bağımsız doğrulama\n• Sunum provasında izleyici gözüyle geri bildirim\n• Final/ödül aşamasında takım temsili desteği', GREEN, 16, 10.8)
card(slide, 9.95, 1.9, 2.85, 3.95, 'Mustafa Gökmen\nDanışman', '• Lise takımı için danışmanlık ve süreç takibi\n• Etik, gizlilik ve şartname uyumluluğu kontrolü\n• Takvim, risk ve resmi iletişim gözetimi\n• Final hazırlığında mentorluk', ORANGE, 16, 10.8)
textbox(slide, 0.85, 6.25, 11.7, 0.45, 'Çalışma prensibi: Ömer liderliğinde hızlı karar alma; diğer roller düşük yoğunluklu ama gerçek, kontrol edilebilir ve jüriye açıklanabilir katkılar üretir.', 13, WHITE, True, PP_ALIGN.CENTER)

# Slide 7 Working model
slide=prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide, 'Çalışma Modeli ve Risk Yönetimi', 'Her aşamada ölçülebilir çıktı, kanıt dosyası ve yedek plan.', '07')
card(slide, 0.75, 1.85, 3.8, 1.6, 'Karar Mekanizması', 'Teknik kararlar Ömer tarafından alınır; danışman şartname/etik uygunluk kontrolü yapar.', CYAN, 16, 11)
card(slide, 4.8, 1.85, 3.8, 1.6, 'Kanıt Odaklı Takip', 'Her kritik adım repo dosyası, rapor, test çıktısı veya submission dosyası ile belgelenir.', GREEN, 16, 11)
card(slide, 8.85, 1.85, 3.8, 1.6, 'Jüriye Hazırlık', 'Teknik derinlik ile sade anlatı ayrılır; canlı demo için 5 dakikalık prova akışı hazırlanır.', ORANGE, 16, 11)
# risk matrix
textbox(slide, 0.8, 4.0, 4.4, 0.4, 'Ana Riskler ve Önlemler', 20, WHITE, True)
risks=[('Veri gecikmesi / format farkı','prepare_kaggle_data.py tek giriş noktası'),('Skor belirsizliği','CV + threshold + ablation ile hızlı iterasyon'),('CPU/GPU limitleri','DistilBERTurk, quantization ve batch optimizasyonu'),('Sunumda karmaşıklık','Basit/teknik modlu anlatı ve XAI ekranı')]
for i,(r,m) in enumerate(risks):
    y=4.55+i*0.48
    pill(slide, 0.85, y, 3.0, 0.34, r, RGBColor(34, 48, 77), WHITE, 9, line=GRAY)
    textbox(slide, 4.05, y+0.04, 7.9, 0.28, '→  ' + m, 10.5, MUTED)

# Slide 8 Closing
slide=prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide, 'Beklenen Katkı', 'Arama kalitesini artıran, açıklanabilir ve uygulanabilir bir e-ticaret relevance sistemi.', '08')
# impact metrics as big numbers
for i,(num,label,col) in enumerate([('3', 'sınıflı alaka kararı', CYAN),('≤250ms', 'hedef p95 servis gecikmesi', GREEN),('1', 'uçtan uca tekrar üretilebilir pipeline', PURPLE),('XAI', 'jüriye açıklanabilir karar', ORANGE)]):
    x=0.85+i*3.08
    shp=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.0), Inches(2.55), Inches(1.75))
    shp.fill.solid(); shp.fill.fore_color.rgb=NAVY2; shp.line.color.rgb=col; shp.line.width=Pt(1.2)
    textbox(slide, x+0.15, 2.25, 2.25, 0.5, num, 24, col, True, PP_ALIGN.CENTER)
    textbox(slide, x+0.2, 2.90, 2.15, 0.5, label, 11, WHITE, True, PP_ALIGN.CENTER)
card(slide, 1.1, 4.45, 11.1, 1.2, 'Kapanış Mesajı', 'Deep-Pipeline; yüksek skor hedefini, şeffaf mühendislik disiplinini ve etkileyici final anlatısını tek bir proje omurgasında birleştirir. Takım yapımızda teknik liderlik net, destek rolleri ise ölçülebilir ve resmi sürece uyumludur.', CYAN, 18, 13)
textbox(slide, 1.0, 6.2, 11.3, 0.45, 'Deep-Pipeline • TEKNOFEST 2026 E-Ticaret Hackathonu', 18, WHITE, True, PP_ALIGN.CENTER)
textbox(slide, 1.0, 6.63, 11.3, 0.25, 'GitHub: github.com/oomerevren/deeppipelinedsad', 10.5, MUTED, False, PP_ALIGN.CENTER)

# Set reading order: bg likely on top? It was first added, so OK; ensure each slide has no placeholders.
OUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(OUT)
print(OUT)
